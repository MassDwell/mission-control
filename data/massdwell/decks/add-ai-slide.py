#!/usr/bin/env python3
"""Add AI Operations slide to v2.3 deck and save as v3.0"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from copy import deepcopy
import os

# Load the converted v2.3 deck
input_path = "/Users/openclaw/Downloads/MassDwell_Investor_Deck_v2.3.pptx"
output_path = "/Users/openclaw/.openclaw/workspace/data/massdwell/decks/MassDwell_Investor_Deck_v3.0.pptx"

prs = Presentation(input_path)

print(f"Loaded presentation with {len(prs.slides)} slides")
print(f"Slide dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")

# Find a good slide to copy style from (slide 9 = Team slide based on the PDF order)
# We'll insert after the Team slide
# Let's examine the slides
for i, slide in enumerate(prs.slides):
    # Get slide title if exists
    title = "No title"
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text[:50]
            if text.strip():
                title = text.strip()[:50]
                break
    print(f"Slide {i+1}: {title}")

print("\n--- Adding AI Operations slide ---")

# Find the Team/Leadership slide (usually after competitive analysis)
# We want to insert after it
# Based on the PDF: Slide 9 is "Leadership"

# Copy the layout from an existing content slide
# We'll duplicate slide 8 (0-indexed = 7) structure
template_idx = 8  # Use slide 9 as template (Leadership)

# Create new slide by duplicating
from pptx.slide import Slide
from lxml import etree

def duplicate_slide(pres, index):
    """Duplicate a slide at the given index"""
    source = pres.slides[index]
    
    # Add blank slide
    blank_layout = pres.slide_layouts[6]  # Blank layout
    new_slide = pres.slides.add_slide(blank_layout)
    
    # Copy shapes from source
    for shape in source.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_slide

# Instead of duplicating, let's add a slide with similar styling
# Add at the end first, then we can reorder in LibreOffice if needed
blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

try:
    new_slide = prs.slides.add_slide(blank_layout)
    print("Added new blank slide")
except Exception as e:
    print(f"Error adding slide: {e}")
    # Try first layout
    new_slide = prs.slides.add_slide(prs.slide_layouts[0])
    print("Added slide with first layout")

# The slide will inherit some styling from the layout
# Add content boxes with the AI Operations content

# Title
title_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "AI-Powered Operations"
p.font.size = Pt(36)
p.font.bold = True
try:
    p.font.color.rgb = RGBColor(96, 165, 250)  # Blue
except:
    pass

# Subtitle
sub_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "We don't just build with AI. We run with AI."
p.font.size = Pt(20)
p.font.italic = True
try:
    p.font.color.rgb = RGBColor(167, 139, 250)  # Purple
except:
    pass

# Content
content_items = [
    ("🤖", "AI Voice Assistant", "Qualifies leads 24/7, books follow-ups automatically"),
    ("📊", "Intelligent CRM", "No lead goes 24+ hours without follow-up"),
    ("📧", "Automated Outreach", "Email triage, drafting, and response scheduling"),
    ("📈", "Real-Time Analytics", "Pipeline forecasting and market monitoring"),
]

y_pos = 2.2
for emoji, title, desc in content_items:
    # Icon/emoji
    icon_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(0.6), Inches(0.5))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(28)
    
    # Title
    title_box = new_slide.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(3), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    
    # Description
    desc_box = new_slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.35), Inches(8), Inches(0.4))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    
    y_pos += 0.95

# Tagline at bottom
tag_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.4))
tf = tag_box.text_frame
p = tf.paragraphs[0]
p.text = "Small team. Massive output. AI amplifies every role."
p.font.size = Pt(16)
p.font.bold = True
try:
    p.font.color.rgb = RGBColor(74, 222, 128)  # Green
except:
    pass

# Save
prs.save(output_path)
print(f"\nSaved: {output_path}")
print(f"Total slides now: {len(prs.slides)}")
