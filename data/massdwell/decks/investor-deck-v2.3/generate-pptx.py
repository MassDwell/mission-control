#!/usr/bin/env python3
"""Generate MassDwell Investor Deck v2.3 PPTX"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors as RGB tuples (will be applied via .rgb property)
def rgb(r, g, b):
    """Return RGB int for pptx"""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

BLUE = rgb(96, 165, 250)
PURPLE = rgb(167, 139, 250)
DARK_BG = rgb(26, 26, 46)
LIGHT_TEXT = rgb(240, 240, 240)
GRAY_TEXT = rgb(136, 136, 136)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = PURPLE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Logo
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(2), Inches(0.4))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MassDwell"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(12)
    
    return slide

# Slide 1: Cover
slide = add_title_slide(prs, "MassDwell", "Where Life Fits")
# Add subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "SAFE INVESTMENT • AI-Powered Modular ADUs"
p.font.size = Pt(18)
p.font.color.rgb = GRAY_TEXT
p.alignment = PP_ALIGN.CENTER

# Slide 2: The Problem
add_content_slide(prs, "America's Housing Crisis", [
    "🏠 3.8M Unit Shortage — National housing deficit creating unprecedented demand",
    "",
    "💰 52% Cost-Burdened — More than half of renters spending over 30% on housing",
    "",
    "👷 40% Labor Shortage — Construction industry facing critical skilled labor deficit",
])

# Slide 3: The Solution  
add_content_slide(prs, "Modular ADUs, Reimagined", [
    "🏭 Factory-Built — Controlled environment for consistent quality, 50% faster",
    "",
    "🤖 AI-Optimized — Generative design from natural language descriptions",
    "",
    "🌱 Sustainable — Net-zero ready construction, eco-friendly materials",
    "",
    "⚡ Turnkey — Design through installation, keys in 8-12 weeks",
])

# Slide 4: Product Line
add_content_slide(prs, "Product Line", [
    "Dwell Essential — 470 sq ft, 1 bed/1 bath — $141,000",
    "",
    "Dwell Classic — 565 sq ft, 2 bed/1 bath — $172,000",
    "",
    "Dwell Deluxe — 594 sq ft, 2 bed/1 bath — $186,000",
    "",
    "Dwell Prime — 892 sq ft, 2 bed/2 bath — $270,000",
    "",
    "Average: ~$192K • Baseline: $300/sq ft turnkey",
])

# Slide 5: Technology
add_content_slide(prs, "AI + Robotics Integration", [
    "AI PLATFORM:",
    "• Generative design from natural language",
    "• Automatic zoning compliance checks",
    "• Predictive demand forecasting",
    "",
    "ROBOTIC MANUFACTURING:",
    "• Sub-millimeter precision (±0.5mm)",
    "• Automated CFS fabrication",
    "• Computer vision QC, zero on-site rework",
])

# Slide 6: Market
add_content_slide(prs, "Market Opportunity", [
    "$47B — Global ADU Market by 2035 (9.19% CAGR)",
    "",
    "MARKET DRIVERS:",
    "• Housing affordability crisis intensifying",
    "• Massachusetts rapidly easing ADU zoning rules",
    "• Multigenerational living on the rise",
    "• Remote work driving flexible space demand",
])

# Slide 7: Competitive Position
add_content_slide(prs, "Competitive Position", [
    "First AI-driven modular manufacturer in the Northeast",
    "",
    "⚡ 1 Month production — vs 6-7 months from competitors",
    "",
    "🎯 ±0.5mm precision — Robotic tolerances no local vendor offers",
    "",
    "🤖 Infinite custom designs — AI-powered vs fixed templates",
    "",
    "West Coast innovators are 3,000+ miles away. Northeast traditionals have no tech.",
])

# Slide 8: Team
add_content_slide(prs, "Leadership", [
    "STEVE VETTORI — Founder & CEO",
    "Real estate developer since 2016. Licensed broker.",
    "Former Atlassian & Cohesity. Principal at Alpine Property Group.",
    "",
    "CARLOS FERREIRA — Founder & CTO", 
    "Technical leadership in manufacturing and construction technology.",
    "Engineering expertise in modular systems and automation.",
    "",
    "Partners: MCSteel (engineered framing) • Alpine Property Group (development)",
])

# Slide 9: Financials
add_content_slide(prs, "Growth Trajectory", [
    "              2026      2027      2028      2029",
    "Units         15        40        100       200",
    "Revenue       $3M       $8M       $20M      $50M",
    "Gross Margin  25%       30%       35%       45%",
    "EBITDA        -$500K    $500K     $4M       $12M",
    "",
    "Break-even: Late 2027 • Target margin: 35-45% at scale",
])

# Slide 10: The Ask
add_content_slide(prs, "The Ask — SAFE Note", [
    "VALUATION CAP: $15M (~5x projected 2026 revenue)",
    "",
    "• Instrument: Post-Money SAFE (YC standard)",
    "• Raising: $1M - $5M",
    "• Minimum investment: $25,000",
    "",
    "YOUR UPSIDE:",
    "• Series A at $30M → 2x advantage",
    "• Series A at $45M → 3x advantage", 
    "• Series A at $75M → 5x advantage",
])

# Slide 11: Use of Funds
add_content_slide(prs, "Use of Funds", [
    "50% — Factory operations & first production units",
    "",
    "25% — AI platform development (Inhabit)",
    "",
    "25% — Team expansion & operations",
    "",
    "2026 MILESTONES:",
    "✓ Factory operational Q1",
    "✓ 15+ units delivered",
    "✓ $3M+ revenue",
    "✓ Path to Series A",
])

# Slide 12: Why SAFE
add_content_slide(prs, "Why SAFE?", [
    "⚡ SPEED — Standard YC docs. Close in days, not months.",
    "",
    "🎯 FOCUS — Defer valuation until we have production data.",
    "",
    "🤝 ALIGNMENT — No board seats or control terms. Execution focus.",
    "",
    "📈 UPSIDE — $15M cap gives early believers significant advantage.",
    "",
    '"We\'re not asking you to value us. We\'re asking you to believe in us."',
])

# Slide 13: Contact
slide = add_title_slide(prs, "MassDwell", "Where Life Fits")
contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(1.5))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "Steve Vettori, Founder & CEO"
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_TEXT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "steve.vettori@massdwell.com • massdwell.com"
p.font.size = Pt(18)
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Needham, Massachusetts"
p.font.size = Pt(16)
p.font.color.rgb = GRAY_TEXT
p.alignment = PP_ALIGN.CENTER

# Save
output_path = "MassDwell_Investor_Deck_v2.3_SAFE.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
